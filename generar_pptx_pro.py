#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
generar_pptx_pro.py  --  Presentacion profesional de ciberseguridad
Requiere: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── PALETA ──────────────────────────────────────────────────────────────────
BG       = RGBColor(0x07, 0x0C, 0x18)   # deep navy
BG_CARD  = RGBColor(0x0F, 0x17, 0x2A)   # card dark
BG_CARD2 = RGBColor(0x19, 0x25, 0x40)   # card mid
BLUE     = RGBColor(0x3B, 0x82, 0xF6)   # primary blue
BLUE_LT  = RGBColor(0x60, 0xA5, 0xFA)   # light blue
BLUE_DIM = RGBColor(0x1D, 0x4E, 0xD8)   # dark blue
GREEN    = RGBColor(0x10, 0xB9, 0x81)   # emerald
RED      = RGBColor(0xF4, 0x3F, 0x5E)   # rose red
AMBER    = RGBColor(0xF5, 0x9E, 0x0B)   # amber
TEAL     = RGBColor(0x06, 0xB6, 0xD4)   # cyan
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GRAY     = RGBColor(0x94, 0xA3, 0xB8)   # slate gray
GRAY_DK  = RGBColor(0x2D, 0x3A, 0x50)   # dark divider
SLATE    = RGBColor(0x64, 0x74, 0x8B)   # muted

W = Inches(13.33)
H = Inches(7.5)

# ── PRESENTACION ─────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]

# ── HELPERS ──────────────────────────────────────────────────────────────────
def bg(slide, color=BG):
    r = slide.shapes.add_shape(1, 0, 0, W, H)
    r.fill.solid()
    r.fill.fore_color.rgb = color
    r.line.fill.background()

def rect(slide, x, y, w, h, fill=None, line_color=None, line_pt=0.75):
    r = slide.shapes.add_shape(1, x, y, w, h)
    if fill:
        r.fill.solid()
        r.fill.fore_color.rgb = fill
    else:
        r.fill.background()
    if line_color:
        r.line.color.rgb = line_color
        r.line.width = Pt(line_pt)
    else:
        r.line.fill.background()
    return r

def txt(slide, text, x, y, w, h, size=14, bold=False, italic=False,
        color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    rn = p.add_run()
    rn.text = text
    rn.font.size = Pt(size)
    rn.font.bold = bold
    rn.font.italic = italic
    rn.font.color.rgb = color
    rn.font.name = "Calibri"
    return tb

def multiline(slide, items, x, y, w, h, def_size=13, def_color=GRAY,
              align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, str):
            d = {'text': item, 'size': def_size, 'bold': False,
                 'color': def_color, 'italic': False}
        else:
            d = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        rn = p.add_run()
        rn.text = d.get('text', '')
        rn.font.size = Pt(d.get('size', def_size))
        rn.font.bold = d.get('bold', False)
        rn.font.italic = d.get('italic', False)
        rn.font.color.rgb = d.get('color', def_color)
        rn.font.name = "Calibri"
    return tb

def title_bar(slide, text, y=Inches(0.28)):
    txt(slide, text, Inches(0.55), y, Inches(12.2), Inches(0.65),
        size=30, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    rect(slide, Inches(0.55), y + Inches(0.72), Inches(12.23), Pt(1.5),
         fill=BLUE)

def subtitle(slide, text, y=Inches(1.15)):
    txt(slide, text, Inches(0.55), y, Inches(12.2), Inches(0.35),
        size=13, color=GRAY, italic=True)

def dot(slide, x, y, color=BLUE, size=Inches(0.08)):
    rect(slide, x, y, size, size, fill=color)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — PORTADA
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
rect(s, 0, 0, W, Inches(0.07), fill=BLUE)

txt(s, "CIBERSEGURIDAD",
    Inches(1.0), Inches(1.7), Inches(11.33), Inches(1.5),
    size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, "EN LA REPUBLICA DOMINICANA",
    Inches(1.0), Inches(3.2), Inches(11.33), Inches(0.7),
    size=26, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

rect(s, Inches(3.5), Inches(4.05), Inches(6.33), Pt(1.5), fill=GRAY_DK)

txt(s, "Amenazas   |   Marco Legal   |   Respuesta Institucional",
    Inches(1.0), Inches(4.2), Inches(11.33), Inches(0.4),
    size=15, color=GRAY, align=PP_ALIGN.CENTER)

txt(s, "Ley No. 53-07  |  DICAT  |  Convenio de Budapest  |  2025",
    Inches(1.0), Inches(6.85), Inches(11.33), Inches(0.35),
    size=11, color=SLATE, align=PP_ALIGN.CENTER)

rect(s, 0, H - Inches(0.07), W, Inches(0.07), fill=BLUE_DIM)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — RESUMEN EJECUTIVO
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
title_bar(s, "RESUMEN EJECUTIVO")
subtitle(s, "Ciberseguridad en la Republica Dominicana — cifras y hallazgos clave")

stats = [
    ("2007", "Promulgacion de la\nLey No. 53-07\nMarco legal vigente", BLUE),
    ("2012", "Primer pais del hemisferio\nen ratificar el\nConvenio de Budapest", GREEN),
    ("30", "Anos de reclusion maxima\npor terrorismo tecnologico\n(Art. 28 — Ley 53-07)", RED),
]
sx = Inches(0.55)
for number, label, color in stats:
    rect(s, sx, Inches(1.55), Inches(3.9), Inches(2.3), fill=BG_CARD,
         line_color=color, line_pt=0.75)
    txt(s, number, sx, Inches(1.65), Inches(3.9), Inches(1.05),
        size=56, bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(s, label, sx + Inches(0.1), Inches(2.7), Inches(3.7), Inches(1.1),
        size=11, color=GRAY, align=PP_ALIGN.CENTER)
    sx += Inches(4.2)

insights = [
    {'text': "Marco Legal Vigente:", 'size': 13, 'bold': True, 'color': BLUE},
    {'text': "La Ley 53-07 tipifica acceso ilicito, fraudes electronicos, robo de identidad y actos de terrorismo tecnologico. Aplica a actores locales y extranjeros con efectos en territorio dominicano.", 'size': 12, 'color': GRAY},
    {'text': " ", 'size': 8, 'color': GRAY},
    {'text': "Posicion Internacional:", 'size': 13, 'bold': True, 'color': GREEN},
    {'text': "La RD es referente regional con cooperacion activa con el Consejo de Europa, la Union Europea, INTERPOL y la Guardia Civil espanola.", 'size': 12, 'color': GRAY},
    {'text': " ", 'size': 8, 'color': GRAY},
    {'text': "Desafio Critico:", 'size': 13, 'bold': True, 'color': AMBER},
    {'text': "La subnotificacion de delitos digitales sigue siendo el principal obstaculo para medir y combatir eficazmente el cibercrimen en el pais.", 'size': 12, 'color': GRAY},
]
multiline(s, insights, Inches(0.55), Inches(4.05), Inches(12.23), Inches(3.2))


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — PANORAMA DE AMENAZAS
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
title_bar(s, "PANORAMA DE AMENAZAS DIGITALES")
subtitle(s, "Vectores de ataque identificados en el entorno dominicano")

threats = [
    ("PHISHING",
     "Correo electronico fraudulento que suplanta identidades para robar credenciales o instalar malware. Vector mas frecuente.",
     RED),
    ("VISHING",
     "Llamadas de ingenieria social. Atacantes se presentan como soporte tecnico, banco o entidades gubernamentales.",
     AMBER),
    ("SMISHING",
     "SMS maliciosos con enlaces fraudulentos o solicitudes urgentes de datos personales y bancarios.",
     AMBER),
    ("MAN-IN-THE-MIDDLE",
     "Interceptacion de trafico en redes inseguras. El atacante se posiciona entre el usuario y el servidor.",
     BLUE),
    ("MEDIA DROPPING",
     "Dispositivos USB abandonados estrategicamente. Al conectarse, ejecutan codigo malicioso de forma silenciosa.",
     TEAL),
]
cx = Inches(0.35)
for title, desc, color in threats:
    rect(s, cx, Inches(1.55), Inches(2.4), Inches(5.55),
         fill=BG_CARD, line_color=color, line_pt=1.0)
    rect(s, cx, Inches(1.55), Inches(2.4), Inches(0.07), fill=color)
    txt(s, title, cx + Inches(0.12), Inches(1.72), Inches(2.18), Inches(0.45),
        size=12, bold=True, color=color)
    txt(s, desc, cx + Inches(0.12), Inches(2.28), Inches(2.18), Inches(4.5),
        size=11, color=GRAY)
    cx += Inches(2.6)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — TRIADA CIA
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
title_bar(s, "PILARES DE LA SEGURIDAD DE LA INFORMACION")
subtitle(s, "Marco CIA — Confidentiality, Integrity, Availability")

pillars = [
    ("C", "CONFIDENCIALIDAD",
     "Solo usuarios y sistemas autorizados pueden acceder a la informacion.",
     ["Controles de acceso por roles", "Cifrado de datos en reposo y transito",
      "Autenticacion multifactor", "Principio de minimo privilegio"],
     BLUE),
    ("I", "INTEGRIDAD",
     "La informacion es exacta, completa y no ha sido alterada por actores no autorizados.",
     ["Firmas digitales y hashes", "Control de versiones",
      "Registros de auditoria inmutables", "Validacion de integridad de archivos"],
     GREEN),
    ("A", "DISPONIBILIDAD",
     "Los sistemas y datos estan accesibles para usuarios autorizados cuando se necesitan.",
     ["Redundancia y failover activo", "Backups con regla 3-2-1",
      "Mitigacion de DDoS", "SLA y monitoreo continuo"],
     AMBER),
]
cx = Inches(0.45)
for letter, name, desc, bullets, color in pillars:
    rect(s, cx, Inches(1.5), Inches(4.1), Inches(5.65),
         fill=BG_CARD, line_color=color, line_pt=0.75)
    txt(s, letter, cx, Inches(1.55), Inches(4.1), Inches(1.3),
        size=72, bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(s, name, cx + Inches(0.15), Inches(2.8), Inches(3.8), Inches(0.42),
        size=14, bold=True, color=WHITE)
    txt(s, desc, cx + Inches(0.15), Inches(3.25), Inches(3.8), Inches(0.85),
        size=11, color=GRAY)
    by = Inches(4.2)
    for b in bullets:
        dot(s, cx + Inches(0.15), by + Inches(0.09), color=color)
        txt(s, b, cx + Inches(0.3), by, Inches(3.6), Inches(0.3),
            size=11, color=GRAY)
        by += Inches(0.33)
    cx += Inches(4.3)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — MODELO DE RIESGO
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
title_bar(s, "MODELO DE RIESGO CIBERNETICO")
subtitle(s, "La relacion entre amenazas, vulnerabilidades e impacto potencial")

components = [
    ("AMENAZA", "Agente o evento capaz de explotar una debilidad: hackers, malware, errores humanos, desastres naturales.", AMBER, True),
    ("+", "", None, False),
    ("VULNERABILIDAD", "Debilidad en sistemas, procesos o personas: software sin parchear, contrasenas debiles, falta de formacion.", RED, True),
    ("=", "", None, False),
    ("RIESGO", "Probabilidad de que una amenaza explote una vulnerabilidad, multiplicada por el impacto potencial en el negocio.", BLUE, True),
]
cx = Inches(0.35)
for label, desc, color, is_card in components:
    if not is_card:
        txt(s, label, cx, Inches(2.6), Inches(0.75), Inches(0.9),
            size=36, bold=True, color=SLATE, align=PP_ALIGN.CENTER)
        cx += Inches(0.85)
    else:
        rect(s, cx, Inches(1.55), Inches(3.3), Inches(4.7),
             fill=BG_CARD, line_color=color, line_pt=1.0)
        rect(s, cx, Inches(1.55), Inches(3.3), Inches(0.07), fill=color)
        txt(s, label, cx + Inches(0.15), Inches(1.72), Inches(3.0), Inches(0.45),
            size=13, bold=True, color=color)
        txt(s, desc, cx + Inches(0.15), Inches(2.3), Inches(3.0), Inches(3.7),
            size=11, color=GRAY)
        cx += Inches(3.6)

rect(s, Inches(0.55), Inches(6.4), Inches(12.23), Inches(0.72),
     fill=BG_CARD2, line_color=BLUE_DIM)
txt(s, "Principio clave:  Sin vulnerabilidad, la amenaza no tiene efecto. Sin amenaza, la vulnerabilidad no genera riesgo. La seguridad efectiva actua sobre ambos vectores de forma simultanea.",
    Inches(0.75), Inches(6.47), Inches(12.0), Inches(0.6),
    size=12, color=BLUE_LT, italic=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — LEY 53-07
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
title_bar(s, "LEY No. 53-07")
subtitle(s, "Crimenes y Delitos de Alta Tecnologia  --  Promulgada el 23 de abril de 2007")

txt(s, "Primera ley dominicana que tipifica y sanciona los delitos informaticos. Protege la confidencialidad, integridad y disponibilidad de los sistemas de informacion del pais.",
    Inches(0.55), Inches(1.52), Inches(12.23), Inches(0.5),
    size=13, color=GRAY)

scope = [
    ("01", "Origen en territorio dominicano",
     "La accion delictiva se origina dentro del pais, independientemente del destino o la victima."),
    ("02", "Origen externo, efectos locales",
     "La accion proviene del extranjero pero produce efectos en la Republica Dominicana."),
    ("03", "Medios locales, efectos externos",
     "Se utilizan medios ubicados en RD para cometer delitos con efectos en el exterior."),
    ("04", "Complicidad desde el pais",
     "Existe participacion o complicidad de actores ubicados en territorio dominicano."),
]
positions = [
    (Inches(0.35), Inches(2.2)),
    (Inches(6.9),  Inches(2.2)),
    (Inches(0.35), Inches(4.1)),
    (Inches(6.9),  Inches(4.1)),
]
for (ix, iy), (num, title, desc) in zip(positions, scope):
    rect(s, ix, iy, Inches(6.2), Inches(1.65), fill=BG_CARD, line_color=GRAY_DK)
    txt(s, num, ix + Inches(0.15), iy + Inches(0.1), Inches(0.55), Inches(0.5),
        size=22, bold=True, color=BLUE)
    txt(s, title, ix + Inches(0.75), iy + Inches(0.1), Inches(5.3), Inches(0.4),
        size=12, bold=True, color=WHITE)
    txt(s, desc, ix + Inches(0.75), iy + Inches(0.52), Inches(5.3), Inches(1.0),
        size=11, color=GRAY)

rect(s, Inches(0.55), Inches(5.92), Inches(12.23), Inches(0.65),
     fill=BG_CARD2, line_color=RED)
txt(s, "Pena maxima:  30 anos de reclusion por actos de terrorismo mediante sistemas tecnologicos (Art. 28)",
    Inches(0.75), Inches(5.99), Inches(12.0), Inches(0.5),
    size=13, bold=True, color=RED)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — CLASIFICACION DE DELITOS
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
title_bar(s, "CLASIFICACION DE DELITOS DIGITALES")
subtitle(s, "Tipificacion bajo la Ley No. 53-07 -- articulado y sanciones")

delitos = [
    ("Acceso ilicito a sistemas informaticos",          "Art. 6",  "3 meses -- 1 ano",    BLUE),
    ("Interceptacion de datos o senales",               "Art. 9",  "1 -- 3 anos",          BLUE),
    ("Clonacion de dispositivos de acceso",             "Art. 5",  "1 -- 10 anos",         AMBER),
    ("Dano o alteracion de datos",                      "Art. 10", "3 meses -- 1 ano",    BLUE),
    ("Sabotaje informatico",                            "Art. 11", "3 meses -- 2 anos",   BLUE),
    ("Estafa electronica",                              "Art. 15", "3 meses -- 7 anos",   RED),
    ("Robo de identidad digital",                       "Art. 17", "3 meses -- 7 anos",   RED),
    ("Chantaje digital",                                "Art. 16", "1 -- 5 anos",          AMBER),
    ("Fraude en telecomunicaciones",                    "Art. 26", "3 meses -- 10 anos",  AMBER),
    ("Espionaje / amenaza a seguridad nacional",        "Art. 27", "15 -- 30 anos",        RED),
    ("Terrorismo mediante sistemas tecnologicos",       "Art. 28", "20 -- 30 anos",        RED),
]

# Header
rect(s, Inches(0.35), Inches(1.5), Inches(12.63), Inches(0.4), fill=BLUE_DIM)
txt(s, "DELITO",          Inches(0.5),  Inches(1.55), Inches(5.5), Inches(0.35),
    size=11, bold=True, color=WHITE)
txt(s, "ARTICULO",       Inches(6.1),  Inches(1.55), Inches(1.8), Inches(0.35),
    size=11, bold=True, color=WHITE)
txt(s, "SANCION",        Inches(8.05), Inches(1.55), Inches(4.7), Inches(0.35),
    size=11, bold=True, color=WHITE)

ry = Inches(1.93)
for i, (delito, art, sancion, color) in enumerate(delitos):
    bg_row = BG_CARD if i % 2 == 0 else BG_CARD2
    rect(s, Inches(0.35), ry, Inches(12.63), Inches(0.46), fill=bg_row)
    dot(s, Inches(0.42), ry + Inches(0.18), color=color, size=Inches(0.1))
    txt(s, delito,  Inches(0.6),  ry + Inches(0.07), Inches(5.4), Inches(0.36),
        size=11, color=WHITE)
    txt(s, art,     Inches(6.1),  ry + Inches(0.07), Inches(1.8), Inches(0.36),
        size=11, bold=True, color=BLUE_LT)
    txt(s, sancion, Inches(8.05), ry + Inches(0.07), Inches(4.7), Inches(0.36),
        size=11, color=color)
    ry += Inches(0.46)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — ARQUITECTURA INSTITUCIONAL
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
title_bar(s, "ARQUITECTURA INSTITUCIONAL DE RESPUESTA")
subtitle(s, "Organismos especializados creados y designados por la Ley 53-07")

institutions = [
    ("DICAT",
     "Depto. de Investigacion de Crimenes y Delitos de Alta Tecnologia",
     "Policia Nacional",
     ["Investiga denuncias ciudadanas y da soporte judicial",
      "Analisis forense digital y preservacion de evidencia",
      "Inteligencia sobre amenazas ciberneticas emergentes",
      "Coordinacion nacional e internacional de casos"],
     BLUE),
    ("DIDI",
     "Division de Investigacion de Delitos Informaticos",
     "Direccion Nacional de Inteligencia (DNI)",
     ["Casos contra la nacion y el Estado dominicano",
      "Espionaje digital y seguridad nacional",
      "Inteligencia cibernetica estrategica",
      "Operaciones encubiertas y contrainteligencia"],
     TEAL),
    ("CICDAT",
     "Comision Interinstitucional contra Crimenes y Delitos de Alta Tecnologia",
     "Presidida por el Procurador General de la Republica",
     ["Coordinacion entre todas las entidades del Estado",
      "Definicion de politica nacional de ciberseguridad",
      "Estadisticas, reportes y rendicion de cuentas",
      "Desarrollo de capacidades institucionales"],
     GREEN),
    ("PEDATEC",
     "Procuraduria Especializada en Delitos de Alta Tecnologia",
     "Procuraduria General de la Republica",
     ["Judicializacion y acusacion publica de casos",
      "Asesoria legal gratuita en delitos digitales",
      "Vinculacion operativa con DICAT y DIDI",
      "Representacion del Estado en procesos penales"],
     AMBER),
]
positions = [
    (Inches(0.35), Inches(1.5)),
    (Inches(6.79), Inches(1.5)),
    (Inches(0.35), Inches(4.35)),
    (Inches(6.79), Inches(4.35)),
]
for (ix, iy), (sigla, nombre, depend, funcs, color) in zip(positions, institutions):
    rect(s, ix, iy, Inches(6.2), Inches(2.65),
         fill=BG_CARD, line_color=color, line_pt=1.0)
    txt(s, sigla, ix + Inches(0.15), iy + Inches(0.1), Inches(2.0), Inches(0.55),
        size=24, bold=True, color=color)
    txt(s, nombre, ix + Inches(0.15), iy + Inches(0.65), Inches(5.9), Inches(0.38),
        size=11, bold=True, color=WHITE)
    txt(s, depend, ix + Inches(0.15), iy + Inches(1.02), Inches(5.9), Inches(0.3),
        size=10, color=GRAY, italic=True)
    fy = iy + Inches(1.38)
    for fn in funcs:
        dot(s, ix + Inches(0.15), fy + Inches(0.09), color=color, size=Inches(0.08))
        txt(s, fn, ix + Inches(0.3), fy, Inches(5.7), Inches(0.28),
            size=10, color=GRAY)
        fy += Inches(0.31)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — POSICION INTERNACIONAL
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
title_bar(s, "POSICIONAMIENTO INTERNACIONAL")
subtitle(s, "La Republica Dominicana como referente regional en ciberseguridad")

rect(s, Inches(0.55), Inches(1.52), Inches(12.23), Inches(0.82),
     fill=BG_CARD2, line_color=GREEN, line_pt=1.0)
txt(s, "Primer pais del hemisferio occidental en adherirse al Convenio de Budapest  (2012)",
    Inches(0.75), Inches(1.62), Inches(12.0), Inches(0.65),
    size=18, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

milestones = [
    ("2003", "Primera investigacion\noficial por delito\ninformatico en RD"),
    ("2007", "Promulgacion Ley 53-07\nCreacion formal\ndel DICAT"),
    ("2008", "Primera cooperacion\ninternacional con la\nGuardia Civil de Espana"),
    ("2012", "Ratificacion del\nConvenio de Budapest\nprimer pais hemisferio"),
    ("2016", "Proyecto Republica Digital\nParticipacion en\nGLACY+ (CE / UE)"),
]
rect(s, Inches(0.55), Inches(3.52), Inches(12.23), Pt(2), fill=BLUE_DIM)
tx = Inches(0.55)
for year, event in milestones:
    txt(s, year, tx, Inches(2.62), Inches(2.3), Inches(0.48),
        size=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    rect(s, tx + Inches(0.9), Inches(3.43), Inches(0.2), Inches(0.2), fill=BLUE)
    txt(s, event, tx, Inches(3.75), Inches(2.3), Inches(1.4),
        size=10, color=GRAY, align=PP_ALIGN.CENTER)
    tx += Inches(2.5)

rect(s, Inches(0.55), Inches(5.55), Inches(12.23), Inches(1.62), fill=BG_CARD)
txt(s, "RED DE COOPERACION ACTIVA",
    Inches(0.75), Inches(5.65), Inches(5.0), Inches(0.38),
    size=12, bold=True, color=BLUE)
partners = [
    ("Consejo de Europa",        "Marco del Convenio de Budapest y programa GLACY+"),
    ("Union Europea",            "Financiamiento y asistencia tecnica en cibercapacidad"),
    ("INTERPOL / Guardia Civil", "Cooperacion operativa e intercambio de inteligencia"),
    ("ITLA / Univ. APEC",        "Formacion tecnica y produccion academica especializada"),
]
px = Inches(0.75)
for partner, desc in partners:
    txt(s, partner, px, Inches(6.12), Inches(2.85), Inches(0.3),
        size=11, bold=True, color=WHITE)
    txt(s, desc,    px, Inches(6.45), Inches(2.85), Inches(0.6),
        size=10, color=GRAY)
    px += Inches(3.1)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CASO: PHISHING BANCARIO
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
title_bar(s, "CASO DE ESTUDIO: FRAUDE ELECTRONICO  --  PHISHING BANCARIO")
subtitle(s, "Analisis de incidente bajo la Ley 53-07 -- Marco legal aplicable")

# Left panel
rect(s, Inches(0.35), Inches(1.5), Inches(7.5), Inches(5.72),
     fill=BG_CARD, line_color=GRAY_DK)
txt(s, "ESCENARIO", Inches(0.55), Inches(1.62), Inches(4.0), Inches(0.34),
    size=11, bold=True, color=BLUE)
scenario = [
    "Pedro recibe un correo aparentemente enviado por su banco,",
    "solicitando verificar sus datos por 'actividad inusual en su cuenta'.",
    "Hace clic en el enlace e ingresa sus credenciales en una pagina",
    "identica al banco pero alojada en un dominio falso.",
    "Horas despues, su cuenta es vaciada mediante transferencias",
    "no autorizadas desde multiples ubicaciones geograficas.",
]
sy = Inches(2.05)
for line in scenario:
    txt(s, line, Inches(0.55), sy, Inches(7.1), Inches(0.3), size=12, color=GRAY)
    sy += Inches(0.32)

txt(s, "VECTOR DE ATAQUE", Inches(0.55), Inches(4.05), Inches(5.0), Inches(0.34),
    size=11, bold=True, color=AMBER)
for v in ["Correo fraudulento (Phishing)",
          "Sitio web clonado con dominio similar",
          "Ingenieria social por urgencia fabricada"]:
    dot(s, Inches(0.55), sy + Inches(0.07), color=AMBER)
    txt(s, v, Inches(0.72), sy, Inches(6.9), Inches(0.3), size=11, color=GRAY)
    sy += Inches(0.32)
    if sy > Inches(4.0):
        sy = Inches(4.5)

txt(s, "IMPACTO EN LA VICTIMA", Inches(0.55), Inches(5.55), Inches(5.5), Inches(0.34),
    size=11, bold=True, color=RED)
txt(s, "Perdida economica total  |  Robo de identidad digital  |  Dano al historial crediticio",
    Inches(0.55), Inches(5.93), Inches(7.1), Inches(0.38), size=11, color=GRAY)

# Right panel
rect(s, Inches(8.05), Inches(1.5), Inches(5.1), Inches(5.72),
     fill=BG_CARD2, line_color=RED, line_pt=1.0)
txt(s, "APLICACION LEGAL", Inches(8.25), Inches(1.62), Inches(4.7), Inches(0.34),
    size=11, bold=True, color=RED)
legal_items = [
    ("Art. 14", "Obtencion ilicita de fondos\nHasta 10 anos de prision"),
    ("Art. 15", "Estafa electronica\n3 meses a 7 anos de prision"),
    ("Art. 8",  "Uso de dispositivos fraudulentos\n1 a 3 anos de prision"),
    ("Art. 17", "Robo de identidad digital\n3 meses a 7 anos de prision"),
]
ly = Inches(2.1)
for art, desc in legal_items:
    rect(s, Inches(8.25), ly, Inches(4.6), Inches(1.12),
         fill=BG_CARD, line_color=GRAY_DK)
    txt(s, art,  Inches(8.4),  ly + Inches(0.08), Inches(1.3), Inches(0.42),
        size=18, bold=True, color=BLUE_LT)
    txt(s, desc, Inches(9.75), ly + Inches(0.08), Inches(3.0), Inches(0.92),
        size=10, color=GRAY)
    ly += Inches(1.22)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — CASO REAL: UBER 2022
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
title_bar(s, "CASO REAL: HACKEO A UBER (2022)")
subtitle(s, "Ingenieria social avanzada -- Aplicacion de la Ley 53-07 al contexto internacional")

facts = [
    ("ACTOR MALICIOSO", "Joven de 18 anos sin herramientas avanzadas"),
    ("METODO PRINCIPAL", "Ingenieria social + MFA Fatigue Attack"),
    ("IMPACTO",         "Acceso total a sistemas internos de Uber"),
]
fx = Inches(0.35)
for label, val in facts:
    rect(s, fx, Inches(1.52), Inches(4.1), Inches(0.88),
         fill=BG_CARD2, line_color=BLUE_DIM)
    txt(s, label, fx + Inches(0.15), Inches(1.58), Inches(3.8), Inches(0.32),
        size=10, bold=True, color=BLUE)
    txt(s, val,   fx + Inches(0.15), Inches(1.87), Inches(3.8), Inches(0.38),
        size=12, bold=True, color=WHITE)
    fx += Inches(4.4)

txt(s, "CADENA DE ATAQUE", Inches(0.55), Inches(2.6), Inches(6.0), Inches(0.36),
    size=12, bold=True, color=AMBER)
steps = [
    ("01", "Compra de credenciales corporativas en mercados del dark web"),
    ("02", "Bombardeo de notificaciones MFA hasta lograr fatiga del empleado"),
    ("03", "Suplantacion en WhatsApp como soporte IT interno de la empresa"),
    ("04", "Acceso a VPN corporativa con credenciales validas comprometidas"),
    ("05", "Escalada de privilegios mediante script de PowerShell encontrado en red interna"),
    ("06", "Exfiltracion de datos y comunicacion publica del incidente en Slack"),
]
sy = Inches(3.05)
for num, step in steps:
    rect(s, Inches(0.55), sy + Inches(0.02), Inches(0.4), Inches(0.3), fill=AMBER)
    txt(s, num,  Inches(0.55), sy + Inches(0.02), Inches(0.4), Inches(0.3),
        size=10, bold=True, color=BG, align=PP_ALIGN.CENTER)
    txt(s, step, Inches(1.05), sy, Inches(5.9), Inches(0.32),
        size=11, color=GRAY)
    sy += Inches(0.38)

# Legal analysis
rect(s, Inches(7.3), Inches(2.5), Inches(5.7), Inches(4.72),
     fill=BG_CARD, line_color=AMBER, line_pt=1.0)
txt(s, "TIPIFICACION -- LEY 53-07",
    Inches(7.5), Inches(2.6), Inches(5.3), Inches(0.36),
    size=12, bold=True, color=AMBER)
legal = [
    ("Art. 6",  "Acceso ilicito -- ingreso no autorizado a VPN y servidores"),
    ("Art. 9",  "Interceptacion de datos -- exfiltracion de informacion interna"),
    ("Art. 17", "Robo de identidad -- suplantacion de soporte IT ante empleados"),
    ("Art. 11", "Sabotaje informatico -- comprometer sistemas criticos"),
    ("Art. 16", "Chantaje digital -- amenaza de divulgacion publica de datos"),
]
ly = Inches(3.1)
for art, desc in legal:
    txt(s, art,  Inches(7.5),  ly, Inches(1.05), Inches(0.28),
        size=12, bold=True, color=BLUE_LT)
    txt(s, desc, Inches(8.65), ly, Inches(4.2),  Inches(0.28),
        size=10, color=GRAY)
    ly += Inches(0.5)

rect(s, Inches(7.5), Inches(5.95), Inches(5.3), Inches(1.1),
     fill=BG_CARD2, line_color=GRAY_DK)
txt(s, "Leccion principal:  La ingenieria social supera la mayoria de los controles tecnicos. El factor humano es el vector de entrada mas explotado a nivel global.",
    Inches(7.65), Inches(6.0), Inches(5.0), Inches(1.0),
    size=11, color=BLUE_LT, italic=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — DEEPFAKE DE VOZ
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
title_bar(s, "AMENAZA EMERGENTE: FRAUDE CON DEEPFAKE DE VOZ")
subtitle(s, "Inteligencia artificial como herramienta de ataque -- Nuevo vector de vishing")

rect(s, Inches(0.35), Inches(1.52), Inches(12.63), Inches(0.48),
     fill=BG_CARD2, line_color=RED)
txt(s, "NIVEL DE AMENAZA: CRITICO  --  Tecnologia accesible, alta tasa de exito, deteccion extremadamente dificil",
    Inches(0.55), Inches(1.6), Inches(12.2), Inches(0.36),
    size=12, bold=True, color=RED)

rect(s, Inches(0.35), Inches(2.15), Inches(7.5), Inches(4.05),
     fill=BG_CARD, line_color=GRAY_DK)
txt(s, "ESCENARIO SIMULADO", Inches(0.55), Inches(2.27), Inches(5.5), Inches(0.34),
    size=11, bold=True, color=BLUE)
scenario = [
    "Sofia, contadora de una empresa mediana, recibe una llamada",
    "con la voz exacta de su director financiero, generada por IA.",
    "",
    "La 'voz' ordena realizar una transferencia urgente y confidencial",
    "de RD$500,000 a una cuenta 'temporal' de un proveedor nuevo,",
    "antes del cierre de jornada. La llamada suena completamente real.",
    "",
    "Sofia ejecuta la transferencia. La cuenta pertenecia a los atacantes.",
    "El director financiero real jamas realizo la llamada.",
]
sy = Inches(2.72)
for line in scenario:
    txt(s, line, Inches(0.55), sy, Inches(7.1), Inches(0.3), size=11, color=GRAY)
    sy += Inches(0.28)

rect(s, Inches(8.05), Inches(2.15), Inches(5.1), Inches(4.05),
     fill=BG_CARD, line_color=RED, line_pt=1.0)
txt(s, "COMO FUNCIONA EL ATAQUE",
    Inches(8.25), Inches(2.27), Inches(4.8), Inches(0.34),
    size=11, bold=True, color=RED)
how_items = [
    "1. Recopilacion de audio (redes sociales, videos, entrevistas)",
    "2. Entrenamiento del modelo con < 30 segundos de voz real",
    "3. Sintesis de voz en tiempo real o grabacion previa del guion",
    "4. Llamada al objetivo en momento de alta presion laboral",
]
hy = Inches(2.72)
for h in how_items:
    txt(s, h, Inches(8.25), hy, Inches(4.7), Inches(0.36), size=10, color=GRAY)
    hy += Inches(0.42)

txt(s, "SENALES DE ALERTA",
    Inches(8.25), Inches(4.37), Inches(4.8), Inches(0.34),
    size=11, bold=True, color=AMBER)
alerts = ["Urgencia extrema e inusual en la solicitud",
          "Confidencialidad total requerida por el llamante",
          "Metodo de pago inusual o cuenta nueva de destino",
          "No es posible verificar por otro canal independiente"]
ay = Inches(4.82)
for a in alerts:
    dot(s, Inches(8.25), ay + Inches(0.08), color=AMBER)
    txt(s, a, Inches(8.42), ay, Inches(4.5), Inches(0.3), size=10, color=GRAY)
    ay += Inches(0.33)

rect(s, Inches(0.35), Inches(6.37), Inches(12.63), Inches(0.75),
     fill=BG_CARD2, line_color=GRAY_DK)
txt(s, "Marco Legal:  Art. 15 (Estafa electronica)  |  Art. 14 (Fraude financiero)  |  Art. 17 (Suplantacion de identidad)  |  Art. 19 (Privacidad)",
    Inches(0.55), Inches(6.48), Inches(12.2), Inches(0.55),
    size=12, color=BLUE_LT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — INDICADORES DE COMPROMISO
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
title_bar(s, "INDICADORES DE COMPROMISO Y DETECCION")
subtitle(s, "Senales de alerta para ciudadanos y organizaciones ante ataques digitales")

categories = [
    ("CORREO ELECTRONICO", RED, [
        "Remitente con dominio similar pero incorrecto (paypa1.com)",
        "Urgencia fabricada: 'Su cuenta sera cerrada en 24 horas'",
        "Solicitud de credenciales o datos bancarios por email",
        "Hipervinculos que no coinciden con el dominio oficial",
        "Errores gramaticales o traduccion automatica deficiente",
    ]),
    ("SITIOS WEB / ENLACES", AMBER, [
        "URL sin HTTPS o con certificado de seguridad no verificado",
        "Dominio ligeramente modificado del sitio original",
        "Diseno identico al oficial pero con pequenas diferencias",
        "Formularios que solicitan datos excesivos o innecesarios",
        "Sin pagina 'Sobre nosotros', politica de privacidad o contacto",
    ]),
    ("LLAMADAS Y MENSAJES", TEAL, [
        "Llamada no solicitada pidiendo datos para 'verificacion'",
        "SMS con enlace acortado sobre premios, alertas o deudas",
        "Presion para actuar de inmediato sin consultar a nadie",
        "Solicitud de acceso remoto al dispositivo personal o corporativo",
        "Voz, redaccion o numero que no coincide con el emisor esperado",
    ]),
]
cx = Inches(0.35)
for cat, color, items in categories:
    rect(s, cx, Inches(1.52), Inches(4.1), Inches(5.65),
         fill=BG_CARD, line_color=color, line_pt=1.0)
    rect(s, cx, Inches(1.52), Inches(4.1), Inches(0.07), fill=color)
    txt(s, cat, cx + Inches(0.15), Inches(1.7), Inches(3.8), Inches(0.38),
        size=12, bold=True, color=color)
    iy = Inches(2.18)
    for item in items:
        dot(s, cx + Inches(0.15), iy + Inches(0.09), color=color, size=Inches(0.08))
        txt(s, item, cx + Inches(0.3), iy, Inches(3.65), Inches(0.4),
            size=10, color=GRAY)
        iy += Inches(0.44)
    cx += Inches(4.3)

rect(s, Inches(0.55), Inches(7.2), Inches(12.23), Pt(1.5), fill=GRAY_DK)
txt(s, "Regla de oro:  En caso de duda -- NO hacer clic, NO responder, NO transferir. Verificar siempre por un canal independiente.",
    Inches(0.55), Inches(7.22), Inches(12.23), Inches(0.25),
    size=11, color=BLUE_LT, italic=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — PROTOCOLO DE DENUNCIA
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
title_bar(s, "PROTOCOLO DE DENUNCIA Y RESPUESTA A INCIDENTES")
subtitle(s, "Pasos ante un incidente de ciberseguridad en la Republica Dominicana")

steps = [
    ("01", "PRESERVAR EVIDENCIA",
     "No borrar mensajes, correos ni registros. Tomar capturas de pantalla con fecha y hora visible. No formatear dispositivos comprometidos.",
     BLUE),
    ("02", "DOCUMENTAR EL INCIDENTE",
     "Registrar: fecha, hora, descripcion detallada del evento, datos del posible atacante, dano economico o personal estimado.",
     TEAL),
    ("03", "DENUNCIAR ANTE EL DICAT",
     "Portal digital de la Policia Nacional (denuncias virtuales o anonimas), presencialmente ante el DICAT o en cualquier fiscalia del pais.",
     GREEN),
    ("04", "SOPORTE Y SEGUIMIENTO",
     "Sistema 311 para quejas administrativas. Contactar PEDATEC para orientacion legal especializada y gratuita durante todo el proceso.",
     AMBER),
]
sx = Inches(0.35)
for num, title, desc, color in steps:
    rect(s, sx, Inches(1.52), Inches(3.1), Inches(5.68),
         fill=BG_CARD, line_color=color, line_pt=1.0)
    rect(s, sx, Inches(1.52), Inches(3.1), Inches(1.08), fill=color)
    txt(s, num, sx, Inches(1.57), Inches(3.1), Inches(0.95),
        size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, title, sx + Inches(0.15), Inches(2.72), Inches(2.8), Inches(0.44),
        size=13, bold=True, color=color)
    txt(s, desc,  sx + Inches(0.15), Inches(3.22), Inches(2.8), Inches(3.7),
        size=11, color=GRAY)
    sx += Inches(3.3)

rect(s, Inches(0.55), Inches(7.18), Inches(12.23), Inches(0.28), fill=BG_CARD2)
txt(s, "DICAT -- Policia Nacional  |  PEDATEC -- Procuraduria General  |  311 -- Quejas Administrativas  |  denuncias.policia.gob.do",
    Inches(0.75), Inches(7.2), Inches(12.0), Inches(0.24),
    size=10, color=GRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — RECOMENDACIONES CLAVE
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
title_bar(s, "RECOMENDACIONES CLAVE")
subtitle(s, "Mejores practicas para individuos y organizaciones en el entorno digital")

rect(s, Inches(0.35), Inches(1.52), Inches(6.1), Inches(5.72),
     fill=BG_CARD, line_color=BLUE, line_pt=1.0)
txt(s, "CIUDADANO DIGITAL",
    Inches(0.55), Inches(1.63), Inches(5.7), Inches(0.42),
    size=14, bold=True, color=BLUE)
ind_recs = [
    ("Contrasenas robustas",
     "Minimo 12 caracteres combinando letras, numeros y simbolos. Una contrasena unica por cuenta."),
    ("Autenticacion multifactor",
     "Activar MFA en todas las cuentas criticas: correo, banca en linea y redes sociales."),
    ("Verificar antes de actuar",
     "Antes de hacer clic en un enlace, verificar el remitente real y el dominio completo de la URL."),
    ("Actualizaciones al dia",
     "Mantener el sistema operativo, navegador, antivirus y aplicaciones siempre actualizados."),
    ("Evitar Wi-Fi publico sin VPN",
     "Nunca realizar transacciones bancarias o acceder a cuentas en redes abiertas sin VPN activa."),
    ("Reportar siempre",
     "Un delito no reportado no puede investigarse ni sancionarse. Denunciar ante el DICAT."),
]
ry = Inches(2.15)
for title, desc in ind_recs:
    dot(s, Inches(0.55), ry + Inches(0.1), color=BLUE)
    txt(s, title, Inches(0.72), ry, Inches(5.55), Inches(0.28),
        size=11, bold=True, color=WHITE)
    txt(s, desc,  Inches(0.72), ry + Inches(0.28), Inches(5.55), Inches(0.35),
        size=10, color=GRAY)
    ry += Inches(0.74)

rect(s, Inches(6.68), Inches(1.52), Inches(6.3), Inches(5.72),
     fill=BG_CARD, line_color=GREEN, line_pt=1.0)
txt(s, "ORGANIZACIONES",
    Inches(6.88), Inches(1.63), Inches(5.9), Inches(0.42),
    size=14, bold=True, color=GREEN)
org_recs = [
    ("Politica de seguridad documentada",
     "Definir, comunicar y hacer cumplir la politica de seguridad de la informacion en toda la organizacion."),
    ("Formacion continua del personal",
     "Simulacros de phishing, capacitaciones trimestrales y evaluaciones de ingenieria social."),
    ("Control de acceso granular",
     "Principio de minimo privilegio. Revocar accesos de inmediato al desvincularse cualquier colaborador."),
    ("Backups y plan de recuperacion",
     "Regla 3-2-1: 3 copias, 2 medios distintos, 1 fuera de sitio. Probar la restauracion periodicamente."),
    ("Gestion de vulnerabilidades",
     "Escaneos periodicos y parcheo en menos de 30 dias para vulnerabilidades criticas identificadas."),
    ("Plan de respuesta a incidentes",
     "Procedimiento documentado, roles claramente definidos y simulacros anuales. Reportar al DICAT."),
]
ry = Inches(2.15)
for title, desc in org_recs:
    dot(s, Inches(6.88), ry + Inches(0.1), color=GREEN)
    txt(s, title, Inches(7.05), ry, Inches(5.75), Inches(0.28),
        size=11, bold=True, color=WHITE)
    txt(s, desc,  Inches(7.05), ry + Inches(0.28), Inches(5.75), Inches(0.35),
        size=10, color=GRAY)
    ry += Inches(0.74)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — CONCLUSION
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
rect(s, 0, 0, W, Inches(0.07), fill=BLUE)

txt(s, "La ciberseguridad no es un producto.",
    Inches(1.0), Inches(1.55), Inches(11.33), Inches(0.9),
    size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Es un proceso continuo.",
    Inches(1.0), Inches(2.4), Inches(11.33), Inches(0.9),
    size=44, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

rect(s, Inches(3.5), Inches(3.42), Inches(6.33), Pt(1.5), fill=GRAY_DK)

txt(s, "PUNTOS CLAVE",
    Inches(1.0), Inches(3.62), Inches(11.33), Inches(0.34),
    size=12, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

takeaways = [
    "El riesgo cibernetico es gestionable: reducir vulnerabilidades reduce el impacto de cualquier amenaza.",
    "La Ley 53-07 provee un marco solido -- la subnotificacion ciudadana es el principal obstaculo a superar.",
    "La ingenieria social supera la mayoria de controles tecnicos -- la formacion humana es el escudo critico.",
    "La RD es referente regional: Convenio de Budapest, DICAT, PEDATEC y cooperacion internacional activa.",
    "El ciudadano digital tiene responsabilidad activa: reportar, prevenir y educarse es parte de la solucion.",
]
ty = Inches(4.08)
for t in takeaways:
    dot(s, Inches(1.2), ty + Inches(0.1), color=BLUE)
    txt(s, t, Inches(1.42), ty, Inches(10.5), Inches(0.3), size=12, color=GRAY)
    ty += Inches(0.4)

txt(s, "Ley No. 53-07  |  DICAT  |  Convenio de Budapest  |  GLACY+  |  2025",
    Inches(1.0), Inches(7.1), Inches(11.33), Inches(0.28),
    size=11, color=SLATE, align=PP_ALIGN.CENTER)
rect(s, 0, H - Inches(0.07), W, Inches(0.07), fill=BLUE_DIM)


# ── GUARDAR ──────────────────────────────────────────────────────────────────
output = r"C:\Users\HP\Desktop\cyber secirity thing\Ciberseguridad_PRO.pptx"
prs.save(output)
print("Presentacion profesional guardada en: " + output)
print("Total de slides: " + str(len(prs.slides)))
