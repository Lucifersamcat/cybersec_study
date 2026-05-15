from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

W = 13.33
H = 7.5

# ── color palette ──────────────────────────────
DARK_BG   = RGBColor(0x0D, 0x14, 0x20)
NAVY_BG   = RGBColor(0x11, 0x18, 0x27)
LIGHT_BG  = RGBColor(0xF8, 0xFA, 0xFC)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BLUE      = RGBColor(0x25, 0x63, 0xEB)
LIGHT_BLU = RGBColor(0x4F, 0x8E, 0xF7)
CYAN      = RGBColor(0x08, 0x91, 0xB2)
GREEN     = RGBColor(0x05, 0x96, 0x69)
PURPLE    = RGBColor(0x7C, 0x3A, 0xED)
RED       = RGBColor(0xDC, 0x26, 0x26)
ORANGE    = RGBColor(0xD9, 0x77, 0x06)
DARK_TXT  = RGBColor(0x1E, 0x29, 0x3B)
MUTED     = RGBColor(0x64, 0x74, 0x8B)
BORDER    = RGBColor(0xE2, 0xE8, 0xF0)
CARD_BG   = RGBColor(0xF1, 0xF5, 0xF9)
ICE_BLUE  = RGBColor(0xEF, 0xF6, 0xFF)
YELLOW    = RGBColor(0xF9, 0xC7, 0x4F)

BLANK = prs.slide_layouts[6]  # blank layout

# ── helpers ────────────────────────────────────
def add_rect(slide, x, y, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background() if line is None else None
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, size=12, color=None, bold=False, italic=False,
             align=PP_ALIGN.LEFT, valign=None, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txBox

def slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def section_header(title, subtitle, label):
    slide = prs.slides.add_slide(BLANK)
    slide_bg(slide, NAVY_BG)
    add_rect(slide, 0, 0, 0.4, H, fill=BLUE)
    add_rect(slide, 0.7, 2.2, 2.5, 0.42, fill=BLUE)
    add_text(slide, label, 0.7, 2.2, 2.5, 0.42, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, title, 0.7, 2.8, 11, 1.3, size=42, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.7, 4.3, 11, 0.6, size=16, color=RGBColor(0x94, 0xA3, 0xB8))
    return slide

def header_bar(slide, title, color=DARK_BG, tcolor=WHITE, size=22):
    add_rect(slide, 0, 0, W, 0.9, fill=color)
    add_text(slide, title, 0.5, 0.05, W-1, 0.8, size=size, color=tcolor, bold=True, valign="middle")

def card(slide, x, y, w, h, header_text=None, header_color=BLUE, bg=WHITE):
    add_rect(slide, x, y, w, h, fill=bg, line=BORDER)
    if header_text:
        add_rect(slide, x, y, w, 0.48, fill=header_color)
        add_text(slide, header_text, x+0.05, y+0.02, w-0.1, 0.44,
                 size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════
# SLIDE 1 — PORTADA
# ══════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, DARK_BG)
add_rect(s, 0, 0, 0.38, H, fill=BLUE)
add_rect(s, 0.65, 1.0, 3.5, 0.42, fill=BLUE)
add_text(s, "Ley No. 53-07  •  DICAT  •  Módulos 1 y 2",
         0.65, 1.0, 3.5, 0.42, size=9, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, "Ciberseguridad", 0.65, 1.65, 12, 1.1, size=58, color=WHITE, bold=True)
add_text(s, "y Marco Legal", 0.65, 2.75, 12, 1.0, size=52, color=LIGHT_BLU, bold=True)
add_text(s, "República Dominicana", 0.65, 3.9, 9, 0.55, size=20, color=RGBColor(0x94,0xA3,0xB8))
add_rect(s, 0.65, 4.65, 4.5, 0.05, fill=BLUE)
add_text(s, "Introducción  ·  Ley 53-07  ·  DICAT  ·  Casos Prácticos  ·  Tareas Resueltas",
         0.65, 4.85, 12, 0.4, size=10, color=MUTED)

# ══════════════════════════════════════════════
# SLIDE 2 — ÍNDICE
# ══════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, LIGHT_BG)
header_bar(s, "Contenido de la Presentación")

items = [
    ("01", "Introducción a la Ciberseguridad", "Definición, pilares y conceptos clave"),
    ("02", "Taller Práctico Resuelto", "3 escenarios de amenazas y vulnerabilidades"),
    ("03", "Glosario de Ciberseguridad", "15 términos esenciales definidos"),
    ("04", "Ley No. 53-07", "Delitos, penas e instituciones creadas"),
    ("05", "Módulo 2 — Prácticas", "Caso del profesor, Uber y Deepfake de voz"),
    ("06", "Casos Prácticos", "5 escenarios con preguntas resueltas"),
]
for i, (num, label, desc) in enumerate(items):
    col = 0 if i < 3 else 1
    row = i % 3
    x = 0.5 if col == 0 else 7.0
    y = 1.1 + row * 1.9
    add_rect(s, x, y, 6.0, 1.65, fill=WHITE, line=BORDER)
    add_rect(s, x, y, 0.65, 1.65, fill=BLUE)
    add_text(s, num, x, y, 0.65, 1.65, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, label, x+0.75, y+0.15, 5.1, 0.5, size=14, color=DARK_TXT, bold=True)
    add_text(s, desc, x+0.75, y+0.72, 5.1, 0.75, size=11, color=MUTED)

# ══════════════════════════════════════════════
# SLIDE 3 — SECCIÓN: MÓDULO 1
# ══════════════════════════════════════════════
section_header("Introducción a la Ciberseguridad", "Conceptos fundamentales y amenazas digitales", "MÓDULO 1")

# ══════════════════════════════════════════════
# SLIDE 4 — ¿QUÉ ES LA CIBERSEGURIDAD?
# ══════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, LIGHT_BG)
header_bar(s, "¿Qué es la Ciberseguridad?")

add_rect(s, 0.4, 1.0, 12.5, 1.6, fill=BLUE)
add_text(s, '"El cuidado y protección de la información, propia y de terceros, gestionada en dispositivos electrónicos como computadoras, redes y teléfonos."',
         0.6, 1.1, 12.1, 1.4, size=14, color=WHITE, italic=True, align=PP_ALIGN.CENTER)

pilares = [
    ("🔒", "CONFIDENCIALIDAD", "Solo acceden quienes están autorizados. La información privada permanece privada."),
    ("⚡", "DISPONIBILIDAD", "Garantizar el acceso a la información cuando se necesite, sin interrupciones."),
    ("✅", "INTEGRIDAD", "La información es exacta, completa y no ha sido alterada sin autorización."),
]
for i, (icon, title, desc) in enumerate(pilares):
    x = 0.4 + i * 4.3
    add_rect(s, x, 2.85, 4.1, 4.2, fill=WHITE, line=BORDER)
    add_rect(s, x, 2.85, 4.1, 0.5, fill=BLUE)
    add_text(s, title, x+0.1, 2.88, 3.9, 0.45, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, icon, x, 3.5, 4.1, 0.7, size=28, align=PP_ALIGN.CENTER)
    add_text(s, desc, x+0.2, 4.35, 3.7, 2.5, size=12, color=MUTED, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════
# SLIDE 5 — AMENAZAS VS VULNERABILIDADES
# ══════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, LIGHT_BG)
header_bar(s, "Amenazas vs. Vulnerabilidades")

add_rect(s, 0.4, 1.0, 5.8, 5.9, fill=WHITE, line=BORDER)
add_rect(s, 0.4, 1.0, 5.8, 0.52, fill=ORANGE)
add_text(s, "⚠  VULNERABILIDAD", 0.4, 1.0, 5.8, 0.52, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, "Debilidad o fallo en un sistema que puede ser explotada:", 0.6, 1.65, 5.4, 0.55, size=12, color=DARK_TXT, bold=True)
vuln_items = ["• Contraseñas débiles o reutilizadas","• Software sin actualizar","• Mala configuración del sistema","• Falta de cifrado en comunicaciones","• Errores de diseño en aplicaciones"]
for i, item in enumerate(vuln_items):
    add_text(s, item, 0.6, 2.3+i*0.7, 5.4, 0.6, size=12, color=MUTED)

add_rect(s, 7.1, 1.0, 5.8, 5.9, fill=WHITE, line=BORDER)
add_rect(s, 7.1, 1.0, 5.8, 0.52, fill=RED)
add_text(s, "⚡  AMENAZA", 7.1, 1.0, 5.8, 0.52, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, "Agente o evento capaz de explotar una vulnerabilidad:", 7.3, 1.65, 5.4, 0.55, size=12, color=DARK_TXT, bold=True)
amn_items = ["• Hackers y cibercriminales","• Malware y ransomware","• Errores humanos","• Phishing / Vishing / SMShing","• Desastres naturales"]
for i, item in enumerate(amn_items):
    add_text(s, item, 7.3, 2.3+i*0.7, 5.4, 0.6, size=12, color=MUTED)

add_rect(s, 5.9, 3.35, 1.5, 0.6, fill=RED)
add_text(s, "= RIESGO", 5.9, 3.35, 1.5, 0.6, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════
# SLIDE 6 — SECCIÓN: TALLER
# ══════════════════════════════════════════════
section_header("Taller Práctico", "Escenarios de amenazas — resueltos", "TALLER PRÁCTICO")

# ══════════════════════════════════════════════
# SLIDE 7 — ESCENARIO 1: WI-FI
# ══════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, LIGHT_BG)
header_bar(s, "Escenario 1 — Wi-Fi pública sin contraseña", size=20)

add_rect(s, 0.4, 1.0, 12.5, 0.85, fill=ICE_BLUE, line=RGBColor(0xBF,0xDB,0xFE))
add_text(s, "📋 Juan se conecta a una red Wi-Fi gratuita en una cafetería para revisar su correo. Un atacante en la misma red intercepta su tráfico.", 0.6, 1.05, 12.0, 0.75, size=12, color=DARK_TXT)

add_rect(s, 0.4, 2.05, 6.0, 4.8, fill=WHITE, line=BORDER)
add_rect(s, 0.4, 2.05, 6.0, 0.5, fill=ORANGE)
add_text(s, "⚠  VULNERABILIDAD", 0.4, 2.05, 6.0, 0.5, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, "La red Wi-Fi carece de cifrado y autenticación. Cualquier persona conectada puede observar el tráfico de datos que circula por ella sin restricción alguna.", 0.6, 2.65, 5.7, 4.0, size=12, color=MUTED)

add_rect(s, 6.9, 2.05, 6.0, 4.8, fill=WHITE, line=BORDER)
add_rect(s, 6.9, 2.05, 6.0, 0.5, fill=RED)
add_text(s, "⚡  AMENAZA — MAN-IN-THE-MIDDLE", 6.9, 2.05, 6.0, 0.5, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, "El atacante se posiciona entre Juan y el punto de acceso, interceptando correos, contraseñas y cualquier información transmitida sin cifrar.\n\nPrevención: usar VPN, evitar Wi-Fi públicas para datos sensibles, preferir redes con contraseña.", 7.1, 2.65, 5.7, 4.0, size=12, color=MUTED)

# ══════════════════════════════════════════════
# SLIDE 8 — ESCENARIO 2: CONTRASEÑA DÉBIL
# ══════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, LIGHT_BG)
header_bar(s, 'Escenario 2 — Contraseña débil "12345"', size=20)

add_rect(s, 0.4, 1.0, 12.5, 0.75, fill=ICE_BLUE, line=RGBColor(0xBF,0xDB,0xFE))
add_text(s, "Laura usa '12345' para todas sus cuentas. Un atacante la adivina y accede a sus perfiles.", 0.6, 1.05, 12.0, 0.65, size=12, color=DARK_TXT)

add_rect(s, 0.4, 1.95, 6.0, 2.35, fill=WHITE, line=BORDER)
add_rect(s, 0.4, 1.95, 6.0, 0.5, fill=ORANGE)
add_text(s, "⚠  VULNERABILIDADES", 0.4, 1.95, 6.0, 0.5, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
vuln = ["• Contraseña débil y predecible", "• Reutilización en múltiples cuentas", "• Sin autenticación adicional (2FA)"]
for i, v in enumerate(vuln):
    add_text(s, v, 0.6, 2.55+i*0.55, 5.7, 0.5, size=12, color=MUTED)

add_rect(s, 6.9, 1.95, 6.0, 2.35, fill=WHITE, line=BORDER)
add_rect(s, 6.9, 1.95, 6.0, 0.5, fill=RED)
add_text(s, "⚡  AMENAZA", 6.9, 1.95, 6.0, 0.5, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, "Fuerza bruta y credential stuffing — uso de contraseñas filtradas en otras plataformas para acceder simultáneamente a todas las cuentas.", 7.1, 2.55, 5.7, 1.55, size=12, color=MUTED)

add_rect(s, 0.4, 4.5, 12.5, 2.5, fill=WHITE, line=BORDER)
add_rect(s, 0.4, 4.5, 12.5, 0.5, fill=GREEN)
add_text(s, "🛡  MEDIDAS DE PROTECCIÓN", 0.4, 4.5, 12.5, 0.5, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
medidas = [
    "• Contraseñas únicas y largas (+12 caracteres) con letras, números y símbolos",
    "• Activar autenticación de dos factores (2FA) en todas las cuentas",
    "• Usar un gestor de contraseñas para no reutilizar credenciales"
]
for i, m in enumerate(medidas):
    add_text(s, m, 0.6, 5.1+i*0.6, 12.1, 0.55, size=12, color=MUTED)

# ══════════════════════════════════════════════
# SLIDE 9 — ESCENARIO 3: PHISHING
# ══════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, LIGHT_BG)
header_bar(s, "Escenario 3 — Phishing por correo electrónico", size=20)

add_rect(s, 0.4, 1.0, 12.5, 0.75, fill=ICE_BLUE, line=RGBColor(0xBF,0xDB,0xFE))
add_text(s, "Roberto recibe un correo de su 'banco' pidiendo actualizar datos. El enlace lleva a una página falsa que roba su información.", 0.6, 1.05, 12.0, 0.65, size=12, color=DARK_TXT)

add_rect(s, 0.4, 1.9, 6.0, 2.1, fill=WHITE, line=BORDER)
add_rect(s, 0.4, 1.9, 6.0, 0.5, fill=ORANGE)
add_text(s, "⚠  VULNERABILIDAD", 0.4, 1.9, 6.0, 0.5, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, "Roberto no verifica la autenticidad del correo antes de actuar. Falta de conocimiento sobre señales de phishing.", 0.6, 2.5, 5.7, 1.35, size=12, color=MUTED)

add_rect(s, 6.9, 1.9, 6.0, 2.1, fill=WHITE, line=BORDER)
add_rect(s, 6.9, 1.9, 6.0, 0.5, fill=RED)
add_text(s, "⚡  AMENAZA — PHISHING", 6.9, 1.9, 6.0, 0.5, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, "Suplantación de identidad del banco para robar credenciales financieras mediante una página web clonada idéntica a la original.", 7.1, 2.5, 5.7, 1.35, size=12, color=MUTED)

add_rect(s, 0.4, 4.15, 12.5, 3.0, fill=WHITE, line=BORDER)
add_rect(s, 0.4, 4.15, 12.5, 0.5, fill=GREEN)
add_text(s, "🔍  ¿CÓMO DETECTAR UN CORREO FALSO?", 0.4, 4.15, 12.5, 0.5, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
tips = [
    "• Revisar el dominio del remitente — los bancos no usan Gmail ni dominios genéricos",
    "• No hacer clic en enlaces del correo — ingresar la URL directamente en el navegador",
    "• Los bancos NUNCA piden credenciales por correo electrónico",
    "• Buscar errores ortográficos, urgencia exagerada o diseño genérico en el mensaje",
]
for i, t in enumerate(tips):
    add_text(s, t, 0.6, 4.75+i*0.57, 12.1, 0.5, size=12, color=MUTED)

# ══════════════════════════════════════════════
# SLIDE 10 — GLOSARIO
# ══════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, LIGHT_BG)
header_bar(s, "Glosario de Ciberseguridad", color=GREEN)

terminos = [
    ("Malware", "Software malicioso que daña o infiltra sistemas sin autorización del usuario."),
    ("Ransomware", "Malware que cifra archivos y exige pago para devolver el acceso."),
    ("Firewall", "Sistema que controla el tráfico de red según reglas predefinidas de seguridad."),
    ("2FA", "Autenticación de dos factores: requiere dos formas de verificar la identidad."),
    ("Ingeniería social", "Manipulación psicológica para obtener información confidencial engañando personas."),
    ("Phishing", "Suplantación de entidad confiable vía correo para robar datos sensibles."),
    ("Vishing", "Phishing realizado mediante llamadas telefónicas fraudulentas."),
    ("SMShing", "Phishing realizado mediante mensajes de texto (SMS)."),
    ("Hacking ético", "Ataque autorizado por el propietario para detectar vulnerabilidades."),
    ("Pentesting", "Prueba de penetración para evaluar la seguridad de un sistema."),
    ("VPN", "Red privada virtual que cifra la conexión a internet protegiendo la privacidad."),
    ("Cifrado", "Convierte información en formato ilegible para protegerla de accesos no autorizados."),
    ("Backdoor", "Acceso oculto instalado por un atacante para reingresar sin ser detectado."),
    ("Zero-day", "Vulnerabilidad desconocida explotada antes de que exista un parche de seguridad."),
    ("Media dropping", "Dejar dispositivos (USB) en lugares públicos para que alguien los conecte."),
]
for i, (term, defi) in enumerate(terminos):
    col = 0 if i < 8 else 1
    row = i % 8
    x = 0.35 if col == 0 else 6.85
    y = 1.0 + row * 0.82
    add_rect(s, x, y, 6.2, 0.72, fill=WHITE, line=BORDER)
    add_rect(s, x, y, 0.08, 0.72, fill=GREEN)
    add_text(s, term, x+0.2, y+0.04, 2.0, 0.32, size=11, color=BLUE, bold=True)
    add_text(s, defi, x+0.2, y+0.36, 5.9, 0.32, size=10, color=MUTED)

# ══════════════════════════════════════════════
# SLIDE 11 — VIDEOS INSTRUCTIVOS
# ══════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, LIGHT_BG)
header_bar(s, "Videos Instructivos del Módulo 1")

videos = [
    ("Video 1", "¿Qué es la ciberseguridad?",
     ["• Protección de info en dispositivos electrónicos.",
      "• Tres pilares: Confidencialidad, Disponibilidad e Integridad.",
      "• Trabajo preventivo anticipando amenazas.",
      "• Roles: hacking ético, pentesting, concienciación.",
      "• Conclusión: protege el activo más importante."]),
    ("Video 2", "Amenazas vs. Vulnerabilidades",
     ["• Vulnerabilidad = debilidad en el sistema.",
      "• Amenaza = agente que la explota.",
      "• Sin amenaza la vulnerabilidad no se explota.",
      "• Riesgo = Amenaza + Vulnerabilidad + Impacto.",
      "• Reducir vulnerabilidades baja el riesgo total."]),
    ("Video 3", "La ciberseguridad, tarea de todos",
     ["• Phishing: reportar correos fraudulentos.",
      "• Vishing: identificar llamadas falsas.",
      "• Seguridad física: equipo desbloqueado = riesgo.",
      "• Media dropping: no conectar USB desconocidos.",
      "• La ciberseguridad es responsabilidad de todos."]),
]
for i, (vnum, vtitle, vpuntos) in enumerate(videos):
    x = 0.4 + i * 4.33
    add_rect(s, x, 1.0, 4.1, 6.1, fill=WHITE, line=BORDER)
    add_rect(s, x, 1.0, 4.1, 0.6, fill=BLUE)
    add_text(s, f"▶  {vnum}", x+0.1, 1.0, 4.0, 0.6, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, vtitle, x+0.15, 1.75, 3.8, 0.65, size=13, color=DARK_TXT, bold=True)
    for j, punto in enumerate(vpuntos):
        add_text(s, punto, x+0.15, 2.55+j*0.85, 3.8, 0.75, size=11, color=MUTED)

# ══════════════════════════════════════════════
# SLIDE 12 — SECCIÓN: LEY 53-07
# ══════════════════════════════════════════════
section_header("Ley No. 53-07", "Crímenes y Delitos de Alta Tecnología — 2007", "MARCO LEGAL")

# ══════════════════════════════════════════════
# SLIDE 13 — DELITOS: CONFIDENCIALIDAD
# ══════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, LIGHT_BG)
header_bar(s, "Delitos contra la Confidencialidad e Integridad  (Arts. 5–11)", size=19)

delitos_ci = [
    ("Art. 5", "Códigos de Acceso / Clonación de Dispositivos", "1 – 10 años"),
    ("Art. 6", "Acceso Ilícito a Sistemas Informáticos", "3 meses – 1 año"),
    ("Art. 7", "Acceso para Servicios a Terceros sin Autorización", "3 meses – 1 año"),
    ("Art. 8", "Dispositivos Fraudulentos", "1 – 3 años"),
    ("Art. 9", "Interceptación de Datos o Señales", "1 – 3 años"),
    ("Art. 10", "Daño o Alteración de Datos", "3 meses – 3 años"),
    ("Art. 11", "Sabotaje Informático", "3 meses – 2 años"),
]
for i, (art, nombre, pena) in enumerate(delitos_ci):
    y = 1.05 + i * 0.9
    bg = WHITE if i % 2 == 0 else CARD_BG
    add_rect(s, 0.35, y, 12.6, 0.78, fill=bg, line=BORDER)
    add_rect(s, 0.35, y, 1.15, 0.78, fill=BLUE)
    add_text(s, art, 0.35, y, 1.15, 0.78, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, nombre, 1.65, y+0.12, 8.5, 0.55, size=12, color=DARK_TXT)
    add_rect(s, 10.6, y+0.14, 2.2, 0.5, fill=RGBColor(0xFE,0xF3,0xC7), line=RGBColor(0xFC,0xD3,0x4D))
    add_text(s, pena, 10.6, y+0.14, 2.2, 0.5, size=10, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════
# SLIDE 14 — DELITOS CONTENIDO + NACIÓN
# ══════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, LIGHT_BG)
header_bar(s, "Delitos de Contenido y Contra la Nación  (Arts. 12–28)", size=19)

add_text(s, "Delitos de Contenido (Arts. 12–24)", 0.35, 1.0, 7.5, 0.45, size=13, color=BLUE, bold=True)
contenido = [
    ("Art. 14", "Obtención ilícita de fondos", "3 – 10 años"),
    ("Art. 15", "Estafa electrónica", "3 meses – 7 años"),
    ("Art. 16", "Chantaje digital", "1 – 5 años"),
    ("Art. 17", "Robo de identidad", "3 meses – 7 años"),
    ("Art. 21-22", "Difamación e injuria digital", "3 meses – 1 año"),
    ("Art. 24", "Pornografía infantil", "2 – 4 años"),
]
for i, (art, nombre, pena) in enumerate(contenido):
    y = 1.55 + i * 0.82
    bg = WHITE if i % 2 == 0 else CARD_BG
    add_rect(s, 0.35, y, 8.3, 0.7, fill=bg, line=BORDER)
    add_rect(s, 0.35, y, 1.2, 0.7, fill=PURPLE)
    add_text(s, art, 0.35, y, 1.2, 0.7, size=9, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, nombre, 1.65, y+0.1, 5.0, 0.5, size=11, color=DARK_TXT)
    add_rect(s, 7.0, y+0.12, 1.55, 0.44, fill=RGBColor(0xFE,0xF3,0xC7), line=RGBColor(0xFC,0xD3,0x4D))
    add_text(s, pena, 7.0, y+0.12, 1.55, 0.44, size=9, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

add_text(s, "Contra la Nación y Terrorismo", 9.2, 1.0, 3.8, 0.45, size=13, color=RED, bold=True)
add_rect(s, 9.2, 1.55, 3.8, 2.0, fill=WHITE, line=BORDER)
add_rect(s, 9.2, 1.55, 3.8, 0.5, fill=RED)
add_text(s, "Art. 27 — Sabotaje / Espionaje", 9.2, 1.55, 3.8, 0.5, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, "15 a 30 años\nde reclusión", 9.2, 2.15, 3.8, 0.9, size=18, color=RED, bold=True, align=PP_ALIGN.CENTER)

add_rect(s, 9.2, 3.8, 3.8, 2.0, fill=WHITE, line=BORDER)
add_rect(s, 9.2, 3.8, 3.8, 0.5, fill=RED)
add_text(s, "Art. 28 — Terrorismo Digital", 9.2, 3.8, 3.8, 0.5, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, "20 a 30 años\nde reclusión", 9.2, 4.4, 3.8, 0.9, size=18, color=RED, bold=True, align=PP_ALIGN.CENTER)

add_rect(s, 0.35, 6.7, 12.6, 0.55, fill=RGBColor(0xFE,0xF2,0xF2), line=RGBColor(0xFE,0xCA,0xCA))
add_text(s, "⚠  Personas morales también responden: multas, disolución o clausura (Art. 60)", 0.5, 6.72, 12.3, 0.5, size=10, color=RED)

# ══════════════════════════════════════════════
# SLIDE 15 — INSTITUCIONES
# ══════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, LIGHT_BG)
header_bar(s, "Instituciones Creadas por la Ley 53-07")

inst = [
    ("DICAT", BLUE,   "Depto. de Investigación de Crímenes y\nDelitos de Alta Tecnología", "Policía Nacional", "Investigar y perseguir delitos tecnológicos. Punto de contacto en red internacional 24/7."),
    ("DIDI",  CYAN,   "División de Investigación de\nDelitos Informáticos", "Departamento Nacional de Investigaciones (DNI)", "Casos contra la nación, el Estado y la seguridad nacional."),
    ("CICDAT",PURPLE, "Comisión Interinstitucional contra\nCrímenes y Delitos de Alta Tecnología", "Presidida por el Procurador General", "Coordinación interinstitucional, política nacional y cooperación internacional."),
    ("PEDATEC",GREEN, "Procuraduría Especializada contra\nCrímenes y Delitos de Alta Tecnología", "Ministerio Público", "Documentar y judicializar los casos investigados por el DICAT y la DIDI."),
]
for i, (sigla, color, nombre, dep, fn) in enumerate(inst):
    col = 0 if i < 2 else 1
    row = i % 2
    x = 0.4 if col == 0 else 6.9
    y = 1.05 + row * 3.1
    add_rect(s, x, y, 6.1, 2.7, fill=WHITE, line=BORDER)
    add_rect(s, x, y, 6.1, 0.62, fill=color)
    add_text(s, sigla, x+0.1, y, 6.0, 0.62, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, nombre, x+0.15, y+0.72, 5.8, 0.72, size=11, color=DARK_TXT, bold=True)
    add_text(s, f"Dependencia: {dep}", x+0.15, y+1.48, 5.8, 0.38, size=10, color=BLUE)
    add_text(s, fn, x+0.15, y+1.9, 5.8, 0.65, size=10, color=MUTED)

# ══════════════════════════════════════════════
# SLIDE 16 — HITOS + PROTOCOLO DE DENUNCIA
# ══════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, LIGHT_BG)
header_bar(s, "Hitos Históricos y Protocolo de Denuncia Ciudadana", size=20)

hitos = [
    ("2003", "Primera investigación oficial por delito informático en la República Dominicana."),
    ("2007", "Promulgación de la Ley 53-07 y creación formal del DICAT y la CICDAT."),
    ("2008", "Cooperación con la Guardia Civil de España — primer caso internacional (extorsión)."),
    ("2012", "Ratificación del Convenio de Budapest — RD fue el 1er país del hemisferio."),
    ("2016", "Proyecto República Digital y participación en GLACY+ (Consejo de Europa / UE)."),
]
for i, (año, texto) in enumerate(hitos):
    y = 1.1 + i * 1.15
    add_rect(s, 0.4, y+0.1, 0.8, 0.8, fill=BLUE)
    add_text(s, año, 0.4, y+0.1, 0.8, 0.8, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, texto, 1.4, y+0.12, 5.8, 0.75, size=12, color=DARK_TXT)

add_text(s, "Protocolo de Denuncia", 8.0, 1.0, 5.0, 0.45, size=14, color=GREEN, bold=True)
pasos = [
    ("1", "Preservar evidencia", "No borrar mensajes. Tomar capturas de pantalla."),
    ("2", "Portal Digital", "Denuncia virtual o anónima — Policía Nacional."),
    ("3", "Presencial", "Ante el DICAT o cualquier fiscalía del MP."),
    ("4", "Sistema 311", "Quejas administrativas de delitos tecnológicos."),
]
for i, (n, titulo, desc) in enumerate(pasos):
    y = 1.6 + i * 1.42
    add_rect(s, 8.0, y, 5.1, 1.22, fill=WHITE, line=BORDER)
    add_rect(s, 8.0, y, 0.6, 1.22, fill=GREEN)
    add_text(s, n, 8.0, y, 0.6, 1.22, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, titulo, 8.75, y+0.08, 4.2, 0.42, size=12, color=DARK_TXT, bold=True)
    add_text(s, desc, 8.75, y+0.55, 4.2, 0.55, size=11, color=MUTED)

# ══════════════════════════════════════════════
# SLIDE 17 — SECCIÓN: MÓDULO 2
# ══════════════════════════════════════════════
section_header("Módulo 2 — Prácticas", "Análisis de casos reales y simulados bajo la Ley 53-07", "MÓDULO 2")

# ══════════════════════════════════════════════
# SLIDE 18 — PRÁCTICA 1
# ══════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, LIGHT_BG)
header_bar(s, "Práctica 1 — El correo del profesor", size=20)

add_rect(s, 0.4, 1.0, 12.5, 0.82, fill=ICE_BLUE, line=RGBColor(0xBF,0xDB,0xFE))
add_text(s, "Un estudiante accede sin autorización al correo institucional de un profesor, descarga exámenes, modifica información y comparte los archivos.", 0.6, 1.05, 12.1, 0.72, size=12, color=DARK_TXT)

cols_p1 = [
    ("Delitos (Ley 53-07)", BLUE,
     ["Art. 6 — Acceso ilícito","Art. 6 Párr. I — Revelación de datos","Art. 10 — Alteración de datos","Art. 7 — Acceso para terceros"]),
    ("Consecuencias Legales", RED,
     ["1 a 3 años de prisión","Multa económica","Expulsión de la institución","Cómplices: misma sanción"]),
    ("Medidas de Prevención", GREEN,
     ["Contraseñas únicas + 2FA","Cerrar sesión en equipos compartidos","Capacitación en seguridad digital","Monitoreo de accesos por TI"]),
]
for i, (titulo, color, items) in enumerate(cols_p1):
    x = 0.4 + i * 4.33
    add_rect(s, x, 2.0, 4.1, 5.1, fill=WHITE, line=BORDER)
    add_rect(s, x, 2.0, 4.1, 0.52, fill=color)
    add_text(s, titulo, x+0.1, 2.0, 4.0, 0.52, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        add_text(s, f"• {item}", x+0.15, 2.65+j*1.05, 3.8, 0.9, size=11, color=MUTED)

# ══════════════════════════════════════════════
# SLIDE 19 — PRÁCTICA 2: UBER
# ══════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, LIGHT_BG)
header_bar(s, "Práctica 2 — Caso Real: Hackeo a Uber (2022)", size=20)

add_rect(s, 0.4, 1.0, 12.5, 0.92, fill=ICE_BLUE, line=RGBColor(0xBF,0xDB,0xFE))
add_text(s, "Un joven de 18 años infiltró los sistemas de Uber usando ingeniería social: se hizo pasar por soporte de TI y convenció a un empleado de entregar sus credenciales de acceso.", 0.6, 1.05, 12.1, 0.82, size=12, color=DARK_TXT)

uber_cards = [
    ("Tipo de Delito", BLUE, "Ingeniería social + acceso ilícito a sistemas internos, interceptación de datos confidenciales de empleados y acceso a infraestructura crítica de la empresa."),
    ("Relación Ley 53-07", PURPLE, "Art. 6 (acceso ilícito), Art. 6 Párr. I (revelación de datos), Art. 9 (interceptación de datos pertenecientes a otra entidad)."),
    ("Consecuencias", RED, "Arrestado en el Reino Unido. Declarado culpable en 2023. Enfrentó años de prisión por hackeo y fraude informático bajo la legislación británica."),
    ("Lección + Prevención", GREEN, "Capacitar en ingeniería social. MFA obligatorio. Política de privilegio mínimo. Ninguna tecnología protege si el factor humano falla."),
]
for i, (titulo, color, texto) in enumerate(uber_cards):
    col = 0 if i < 2 else 1
    row = i % 2
    x = 0.4 if col == 0 else 6.9
    y = 2.1 + row * 2.55
    add_rect(s, x, y, 6.1, 2.3, fill=WHITE, line=BORDER)
    add_rect(s, x, y, 0.1, 2.3, fill=color)
    add_text(s, titulo, x+0.25, y+0.1, 5.7, 0.38, size=12, color=color, bold=True)
    add_text(s, texto, x+0.25, y+0.56, 5.7, 1.6, size=11, color=MUTED)

# ══════════════════════════════════════════════
# SLIDE 20 — PRÁCTICA 3: DEEPFAKE DE VOZ
# ══════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, LIGHT_BG)
header_bar(s, "Práctica 3 — Caso Creado: Fraude con Deepfake de Voz", size=19)

add_rect(s, 0.4, 1.0, 12.5, 1.05, fill=ICE_BLUE, line=RGBColor(0xBF,0xDB,0xFE))
add_text(s, "Sofía, contadora, recibe una llamada con la voz clonada de su jefe (IA). Le pide urgentemente transferir RD$800,000. Al hacerlo, el dinero desaparece en cuentas en el extranjero. Los atacantes usaron audios públicos del director para clonar su voz.", 0.6, 1.05, 12.1, 0.95, size=12, color=DARK_TXT)

add_rect(s, 0.4, 2.2, 6.0, 2.3, fill=WHITE, line=BORDER)
add_rect(s, 0.4, 2.2, 6.0, 0.5, fill=BLUE)
add_text(s, "Delitos — Ley 53-07", 0.4, 2.2, 6.0, 0.5, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
for j, (art, txt) in enumerate([("Art. 14","Obtención ilícita de fondos — 3 a 10 años"),("Art. 15","Estafa electrónica — 3 meses a 7 años"),("Art. 17","Robo de identidad — 3 meses a 7 años")]):
    add_text(s, f"• {art}: {txt}", 0.6, 2.82+j*0.55, 5.7, 0.48, size=11, color=MUTED)

add_rect(s, 6.9, 2.2, 6.0, 2.3, fill=WHITE, line=BORDER)
add_rect(s, 6.9, 2.2, 6.0, 0.5, fill=GREEN)
add_text(s, "Plan de Prevención (4 acciones)", 6.9, 2.2, 6.0, 0.5, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
acciones = ["1. Verificación doble para transferencias grandes","2. Capacitación sobre deepfakes de voz","3. Palabra clave secreta entre directivos","4. Múltiples aprobadores para montos altos"]
for j, a in enumerate(acciones):
    add_text(s, a, 7.1, 2.82+j*0.55, 5.6, 0.48, size=11, color=MUTED)

add_rect(s, 0.4, 4.65, 12.5, 2.5, fill=RGBColor(0xFF,0xF7,0xED), line=RGBColor(0xFE,0xD7,0xAA))
add_text(s, "💡  Reflexión Final", 0.6, 4.72, 12.0, 0.42, size=12, color=ORANGE, bold=True)
add_text(s, "Estos delitos son cada vez más comunes porque la tecnología avanza más rápido que la conciencia de las personas. La ciberseguridad no es solo un asunto del departamento de TI — es responsabilidad de todos. Como ciudadano digital debemos cuestionar lo urgente, verificar antes de actuar y mantenernos informados.", 0.6, 5.2, 12.1, 1.8, size=11, color=MUTED)

# ══════════════════════════════════════════════
# SLIDE 21 — SECCIÓN: CASOS PRÁCTICOS
# ══════════════════════════════════════════════
section_header("Casos Prácticos", "5 escenarios con preguntas resueltas", "CASOS PRÁCTICOS")

# ══════════════════════════════════════════════
# SLIDE 22 — CASOS 1 Y 2
# ══════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, LIGHT_BG)
header_bar(s, "Casos 1 y 2 — Privacidad y Fraude Electrónico")

for i, (num, titulo, color, desc, delitos_txt, resp) in enumerate([
    ("1", "Acceso Ilícito — Instagram", PURPLE,
     "Juan accede a la cuenta de Instagram de su ex-pareja usando una contraseña guardada. Publica contenido íntimo y cambia la contraseña.",
     "Art. 6 (Acceso ilícito)  ·  Art. 10 (Daño de datos)  ·  Art. 19 (Invasión de privacidad)",
     "¿Tener la contraseña lo hace legal? NO. Usar credenciales ajenas sin autorización expresa es acceso ilícito sin importar cómo se obtuvo la clave. Pena: 3 meses a 3 años + multa."),
    ("2", "Phishing Bancario", RED,
     "Pedro ingresa sus credenciales en una página web falsa del banco recibida por correo. Horas después su cuenta es vaciada con transacciones no autorizadas.",
     "Art. 15 (Estafa electrónica)  ·  Art. 14 (Obtención ilícita de fondos)  ·  Art. 8 (Dispositivos fraudulentos)",
     "Errores de Pedro: clic en enlace no verificado, no revisó el dominio, ingresó datos sin validar autenticidad. Prevención: nunca ingresar credenciales desde enlaces de correo. Activar 2FA."),
]):
    x = 0.4 if i == 0 else 7.0
    add_rect(s, x, 1.0, 6.0, 6.1, fill=WHITE, line=BORDER)
    add_rect(s, x, 1.0, 6.0, 0.58, fill=color)
    add_text(s, f"Caso {num} — {titulo}", x+0.1, 1.0, 5.9, 0.58, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, desc, x+0.15, 1.68, 5.7, 1.1, size=11, color=DARK_TXT)
    add_text(s, "Delitos aplicables:", x+0.15, 2.88, 5.7, 0.38, size=11, color=color, bold=True)
    add_text(s, delitos_txt, x+0.15, 3.28, 5.7, 0.55, size=10, color=MUTED)
    add_rect(s, x, 3.95, 6.0, 0.05, fill=BORDER)
    add_text(s, "Respuesta clave:", x+0.15, 4.08, 5.7, 0.38, size=11, color=color, bold=True)
    add_text(s, resp, x+0.15, 4.5, 5.7, 2.45, size=10, color=MUTED)

# ══════════════════════════════════════════════
# SLIDE 23 — CASOS 3, 4 Y 5
# ══════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, LIGHT_BG)
header_bar(s, "Casos 3, 4 y 5 — Identidad, Sistema y Privacidad", size=20)

casos = [
    ("3", "Suplantación\ny Difamación", ORANGE,
     "Estudiante crea cuenta falsa con nombre y foto de un profesor, publicando ofensas e información falsa.",
     "Art. 17 (Robo identidad)\nArt. 21 (Difamación)\nArt. 22 (Injuria pública)",
     "Cómo demostrar: IP de registro, historial de dispositivos y metadatos solicitados con orden judicial."),
    ("4", "Sistema Escolar\nHackeado", BLUE,
     "Carlos explota una vulnerabilidad del sistema académico para cambiar calificaciones a cambio de dinero.",
     "Art. 6 (Acceso ilícito)\nArt. 10 (Alteración datos)\nArt. 11 (Sabotaje)",
     "Descubrir la falla NO autoriza explotarla. Debió reportarla al TI. Pagadores son cómplices (Art. 6)."),
    ("5", "Difusión sin\nConsentimiento", PURPLE,
     "La pareja de Ana difunde fotos íntimas en WhatsApp y redes sociales causando daño emocional y social.",
     "Art. 19 (Invasión privacidad)\nArt. 21 (Difamación)\nArt. 23 (Atentado sexual si aplica)",
     "No es 'problema personal' — es un delito. Denuncia inmediata al DICAT y preservar toda evidencia."),
]
for i, (num, titulo, color, desc, delitos_txt, resp) in enumerate(casos):
    x = 0.35 + i * 4.35
    add_rect(s, x, 1.0, 4.1, 6.1, fill=WHITE, line=BORDER)
    add_rect(s, x, 1.0, 4.1, 0.75, fill=color)
    add_text(s, f"Caso {num}", x+0.1, 1.0, 4.0, 0.35, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, titulo, x+0.1, 1.35, 4.0, 0.4, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, desc, x+0.12, 1.85, 3.85, 1.1, size=11, color=DARK_TXT)
    add_text(s, "Delitos:", x+0.12, 3.05, 3.85, 0.3, size=11, color=color, bold=True)
    add_text(s, delitos_txt, x+0.12, 3.38, 3.85, 0.95, size=10, color=MUTED)
    add_rect(s, x, 4.4, 4.1, 0.05, fill=BORDER)
    add_text(s, "Respuesta:", x+0.12, 4.52, 3.85, 0.32, size=11, color=color, bold=True)
    add_text(s, resp, x+0.12, 4.88, 3.85, 2.05, size=10, color=MUTED)

# ══════════════════════════════════════════════
# SLIDE 24 — CONCLUSIÓN
# ══════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, DARK_BG)
add_rect(s, 0, 0, 0.38, H, fill=BLUE)

add_text(s, "Conclusión", 0.65, 0.5, 12, 1.0, size=46, color=WHITE, bold=True)
add_rect(s, 0.65, 1.65, 4.0, 0.07, fill=BLUE)

puntos = [
    ("1", "Amenaza cotidiana real", "Phishing, redes inseguras y contraseñas débiles son las principales puertas de entrada al cibercrimen."),
    ("2", "La Ley 53-07 protege", "Marco legal sólido que sanciona desde el acceso ilícito hasta el terrorismo digital con penas de hasta 30 años."),
    ("3", "Instituciones especializadas", "DICAT, DIDI, CICDAT y PEDATEC garantizan la persecución del cibercrimen en la República Dominicana."),
    ("4", "Rol del ciudadano digital", "Denunciar, verificar antes de actuar y mantenerse informado es responsabilidad de todos los ciudadanos."),
]
for i, (n, titulo, desc) in enumerate(puntos):
    y = 1.9 + i * 1.32
    add_rect(s, 0.65, y+0.1, 0.7, 0.7, fill=BLUE)
    add_text(s, n, 0.65, y+0.1, 0.7, 0.7, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, titulo, 1.55, y+0.08, 11, 0.42, size=15, color=LIGHT_BLU, bold=True)
    add_text(s, desc, 1.55, y+0.52, 11, 0.65, size=12, color=RGBColor(0x94,0xA3,0xB8))

add_rect(s, 0.65, 7.15, 12.0, 0.2, fill=RGBColor(0x1E,0x29,0x3B))
add_text(s, "Ciberseguridad y Marco Legal RD  •  Ley No. 53-07  •  DICAT  •  Módulos 1 y 2",
         0.65, 7.1, 12.0, 0.35, size=9, color=MUTED, align=PP_ALIGN.CENTER)

# ── GUARDAR ─────────────────────────────────────
output = r"C:\Users\HP\Desktop\cyber secirity thing\Ciberseguridad_RD.pptx"
prs.save(output)
print(f"✅ Presentación guardada en: {output}")
print(f"   Total de slides: {len(prs.slides)}")
